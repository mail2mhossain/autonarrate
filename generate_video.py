import os
import time
import shutil
import numpy as np
import edge_tts
from win32com.client import Dispatch, constants, gencache
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
import re
from pydub import AudioSegment
from moviepy import VideoFileClip, AudioFileClip, AudioClip

def ppt_to_video(ppt_path: str,
                 video_path: str,
                 use_timings: bool = False,
                 default_slide_duration: int = 5,
                 fps: int = 30,
                 vert_resolution: int = 720,
                 quality: int = 80,
                 progress_callback=None):
    """
    Exports a PowerPoint presentation to a video file, using either default slide durations or custom timings/narrations.
    """
    # Launch PowerPoint (headless)
    ppt = Dispatch("PowerPoint.Application")
    # ppt.Visible = False

    # Open presentation
    if progress_callback:
        progress_callback(10, "Opening PowerPoint presentation...")
    pres = ppt.Presentations.Open(os.path.abspath(ppt_path),
                                  WithWindow=False)

    # Create video
    if progress_callback:
        progress_callback(20, "Creating video...")
    pres.CreateVideo(FileName=os.path.abspath(video_path),
                    UseTimingsAndNarrations=use_timings,
                    # DefaultSlideDuration=default_slide_duration,
                    FramesPerSecond=fps,
                    VertResolution=vert_resolution,
                    Quality=quality)
    if progress_callback:
        progress_callback(50, "Video generation in progress...")
    
    # Wait until video generation finishes (queued=1 or in-progress=2)
    status = pres.CreateVideoStatus
    progress_status = 1

    while status in (1, 2):  # 1=Queued, 2=InProgress
        if progress_callback:
            progress_callback(50 + progress_status * 10, "Video generation in progress...")
        time.sleep(1)
        status = pres.CreateVideoStatus
        progress_status += 1
    # Inspect final status (3=Done, 4=Failed)
    if status == 3:
        if progress_callback:
            progress_callback(100, "Video generation completed.")
        print(f"✅ Video saved to {video_path}")
    elif status == 4:
        print(f"❌ Video generation failed with status {status}")
    else:
        print(f"⚠️ Unexpected CreateVideoStatus: {status}")

    # Clean up
    pres.Close()
    ppt.Quit()
    time.sleep(5)


def generate_audio_from_points(ppt_path: str, output_dir: str, progress_callback=None):
    """
    Generates audio segments for each bullet point or image cue in every slide of a PowerPoint presentation.
    Text content is converted to speech using TTS, while images may result in silent audio segments for timing.
    """
    
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(ppt_path)
    slide_count = len(prs.slides)

    audio_map = {}
    for slide_idx, slide in enumerate(prs.slides, start=1):
        audio_map[slide_idx] = []
        point_counter = 1
        image_counter = 0
        percent = int(100 * slide_idx / slide_count)
        if progress_callback:
            progress_callback(percent, f"Generating audio {slide_idx}/{slide_count}")

        for shape in slide.shapes:
            print(f"Shape type of slide {slide_idx}: {shape.shape_type}")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_counter += 1
                print(f"Found image {image_counter} on slide {slide_idx}")

            # if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and image_counter == 2:
                # Process the image
                print(f"Generating silence audio [{slide_idx}]")
                duration = 4
                # silent_clip = AudioClip(lambda t: np.array([0.0]), duration=duration)
                def make_silence(t):
                    if np.isscalar(t):
                        return np.zeros((1, 2))  # 1 sample, stereo
                    else:
                        return np.zeros((len(t), 2))  # multiple samples, stereo
                
                fname = os.path.join(output_dir, f"slide_{slide_idx}_point_{point_counter}.mp3")
                if not os.path.exists(fname):
                    silent_clip = AudioClip(make_silence, duration=duration)
                    silent_clip.write_audiofile(fname, fps=44100, codec='libmp3lame')

                print(f"Generated silence audio [{slide_idx}]: {fname}")
                audio_map[slide_idx].append(fname)
                point_counter += 1
                
            if not getattr(shape, "has_text_frame", False):
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                # skip if no alphanumeric content
                if not text or not re.search(r"\w", text):
                    continue
                try:
                    fname = os.path.join(output_dir, f"slide_{slide_idx}_point_{point_counter}.mp3")
                    if not os.path.exists(fname):
                        communicate = edge_tts.Communicate(text, "en-US-ChristopherNeural")
                        communicate.save_sync(fname)

                    print(f"Generated audio: {fname}")
                    audio_map[slide_idx].append(fname)
                    point_counter += 1
                except AssertionError:
                    print(f"⚠️ Skipping invalid TTS text at slide {slide_idx}, point {point_counter}")
                    continue
    return audio_map

def measure_durations(audio_map, progress_callback=None):
    """
    Calculates the duration (in seconds) of each generated audio segment for every slide.
    Returns a mapping of slide indices to lists of durations.
    """
    durations = {}
    total_items = sum(len(files) for files in audio_map.values())
    completed = 0
    for slide_idx, files in audio_map.items():
        durs = []
        for path in files:
            clip = AudioFileClip(path)
            durs.append(clip.duration)
            clip.close()
            completed += 1
            if progress_callback:
                percent = int(100 * completed / total_items)
                progress_callback(percent, f"Measuring duration {completed}/{total_items}")


        durations[slide_idx] = durs
    return durations

def apply_point_timings(pptx_path: str, durations_map: dict, progress_callback=None):
    """
    Applies precise timing to each bullet point or cue in the PowerPoint presentation,
    ensuring that each appears in sync with its corresponding audio.
    """

    # 1) Start PowerPoint in the background
    pp = gencache.EnsureDispatch("PowerPoint.Application")
    
    # 2) Open your presentation
    pres = pp.Presentations.Open(os.path.abspath(pptx_path))

    # 3) Calculate total points for progress
    total_points = sum(len(durations) for durations in durations_map.values())
    completed = 0

    # 4) For each slide that has audio durations defined:
    for slide in pres.Slides:
        idx = slide.SlideIndex
        if idx not in durations_map:
            continue

        point_durs = durations_map[idx]
        seq = slide.TimeLine.MainSequence
        total_points = seq.Count
        completed = 0

        # — First bullet: appear as soon as the slide opens
        print(f"Total sequence of slide {idx} is {seq.Count}. Times: {point_durs}")
        if seq.Count >= 1:
            first = seq.Item(1)
            first.Timing.TriggerType      = constants.msoAnimTriggerWithPrevious
            first.Timing.TriggerDelayTime = point_durs[0]
            completed += 1
            if progress_callback:
                percent = int(100 * completed / total_points)
                progress_callback(percent, f"Applying timing {completed}/{total_points}")

            print(f"Slide {idx} • Point 1: appear {point_durs[0]:.2f}s after previous")

        # — Subsequent bullets: appear after the previous point’s audio
        for i in range(2, seq.Count + 1):
            effect = seq.Item(i)
            prev_dur = point_durs[i-1] if (i-1) < len(point_durs) else 0.0

            effect.Timing.TriggerType      = constants.msoAnimTriggerAfterPrevious
            effect.Timing.TriggerDelayTime = prev_dur
            completed += 1
            if progress_callback:
                percent = int(100 * completed / total_points)
                progress_callback(percent, f"Applying timing {completed}/{total_points}")
            print(f"Slide {idx} • Point {i}: appear {prev_dur:.2f}s after previous")

        # Optional: if you want the slide itself to advance only after
        # all your points have appeared, sum them up:
        total = sum(point_durs)
        # if len(point_durs) == 1:
        #     total +=5.0
        slide.SlideShowTransition.AdvanceOnTime = True
        slide.SlideShowTransition.AdvanceTime   = total
        if progress_callback:
            percent = 100
            progress_callback(percent, f"Timing applied")
        print(f"Slide {idx}: auto-advance after {total:.2f}s total")

    # 4) Save and clean up
    pres.Save()
    time.sleep(5)
    pres.Close()
    pp.Quit()
    
    print("Done: per-point timings applied.")


# --- combine & merge utilities ---
def combine_audio(audio_map, output_audio, progress_callback=None):
    """
    Flattens and concatenates all generated audio clips into a single audio file, preserving the order of slides and points.
    """
    combined = AudioSegment.empty()
    total_items = sum(len(files) for files in audio_map.values())
    completed = 0
    for slide_idx in sorted(audio_map):
        for fpath in audio_map[slide_idx]:
            completed += 1
            if progress_callback:
                percent = int(100 * completed / total_items)
                progress_callback(percent, f"Combining audio {completed}/{total_items}")
            combined += AudioSegment.from_file(fpath)
    combined.export(output_audio, format="mp3")
    if progress_callback:
        progress_callback(100, "Combining audio completed.")
    print(f"🔊 Combined audio saved to {output_audio}")
    return output_audio

def merge_audio_video(video_path, audio_path, output_video="final_video.mp4", progress_callback=None):
    """
    Overlays the combined audio track onto the generated video, producing a final video with synchronized narration.
    """
    if progress_callback:
        progress_callback(10, "Loading video ...")
    video = VideoFileClip(video_path)
    if progress_callback:
        progress_callback(20, "Loading audio...")
    audio = AudioFileClip(audio_path)
    if progress_callback:
        progress_callback(30, "Merging video and audio...")
    final = video.with_audio(audio)
    if progress_callback:
        progress_callback(60, "Saving final video...")
    final.write_videofile(output_video, codec="libx264", audio_codec="aac")
    if progress_callback:
        progress_callback(100, "Video merge completed.")
    print(f"🎬 Merged video saved to {output_video}")
    return output_video


if __name__ == "__main__":
    PPT_FILE   = "Z:\\CCS\\computational_thinking.pptx"
    video_file_name = os.path.splitext(os.path.basename(PPT_FILE))[0]
    parent_dir = os.path.dirname(os.path.abspath(PPT_FILE))
    video_dir = os.path.join(parent_dir, video_file_name)
    
    if not os.path.exists(video_dir):
        os.makedirs(video_dir)

    VIDEO_FILE = os.path.join(video_dir, video_file_name + ".mp4")
    PPT_VIDEO = os.path.join(video_dir, video_file_name + "_ppt.mp4")

    # Generate audio per bullet point
    audio_dir = os.path.join(video_dir, "audio")
    audio_map = generate_audio_from_points(PPT_FILE, audio_dir)
    
    durations_map = measure_durations(audio_map)
    apply_point_timings(PPT_FILE, durations_map)
    if not os.path.exists(PPT_VIDEO):
        ppt_to_video(PPT_FILE, PPT_VIDEO, use_timings=True, default_slide_duration=7)
    
    combined_audio_file = os.path.join(video_dir, "combined_audio.mp3")
    if not os.path.exists(combined_audio_file):
        combine_audio(audio_map, combined_audio_file)

    if not os.path.exists(VIDEO_FILE):
        merge_audio_video(PPT_VIDEO, combined_audio_file,  VIDEO_FILE)

